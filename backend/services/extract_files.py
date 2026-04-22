import os

ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.txt', '.pptx'}

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return _extract_pdf(file_path)
    elif ext in ('.docx', '.doc'):
        return _extract_docx(file_path)
    elif ext == '.txt':
        return _extract_txt(file_path)
    elif ext == '.pptx':
        return _extract_pptx(file_path)
    return ""

def _extract_pdf(path):
    try:
        import fitz
        doc = fitz.open(path)
        return "\n".join(page.get_text() for page in doc)
    except Exception:
        pass
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception:
        return ""

def _extract_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""

def _extract_txt(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        return ""

def _extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""
